# -*- coding: utf-8 -*-
import csv
from collections import defaultdict

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
plt.rcParams['font.family'] = 'Meiryo'

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

COLS = ['sepal.length', 'sepal.width', 'petal.length', 'petal.width']
COL_JP = {
    'sepal.length': 'がく片の長さ',
    'sepal.width': 'がく片の幅',
    'petal.length': '花弁の長さ',
    'petal.width': '花弁の幅',
}
COLORS = {'Setosa': '#4C72B0', 'Versicolor': '#DD8452', 'Virginica': '#55A868'}

# ---------- データ読み込み ----------
rows = []
with open('iris.csv', newline='', encoding='utf-8') as f:
    reader = csv.DictReader(f)
    for r in reader:
        rows.append(r)

by_variety = defaultdict(list)
for r in rows:
    by_variety[r['variety']].append(r)


def stats(values):
    n = len(values)
    mean = sum(values) / n
    mn, mx = min(values), max(values)
    var = sum((v - mean) ** 2 for v in values) / n
    return mean, var ** 0.5, mn, mx


overall_stats = {}
for c in COLS:
    vals = [float(r[c]) for r in rows]
    overall_stats[c] = stats(vals)

variety_means = {}
for v, rs in by_variety.items():
    variety_means[v] = {c: sum(float(r[c]) for r in rs) / len(rs) for c in COLS}

# ---------- グラフ1: 品種ごとの平均値 棒グラフ ----------
fig, ax = plt.subplots(figsize=(9, 4.8))
varieties = ['Setosa', 'Versicolor', 'Virginica']
x = range(len(COLS))
width = 0.25
for i, v in enumerate(varieties):
    means = [variety_means[v][c] for c in COLS]
    ax.bar([p + i * width for p in x], means, width=width, label=v, color=COLORS[v])
ax.set_xticks([p + width for p in x])
ax.set_xticklabels([COL_JP[c] for c in COLS], fontsize=12)
ax.set_ylabel('平均値 (cm)', fontsize=12)
ax.set_title('品種ごとの特徴量の平均値', fontsize=14)
ax.legend()
fig.tight_layout()
fig.savefig('chart_bar.png', dpi=150)
plt.close(fig)

# ---------- グラフ2: 花弁の長さ×幅 散布図 ----------
fig, ax = plt.subplots(figsize=(9, 4.8))
for v in varieties:
    rs = by_variety[v]
    xs = [float(r['petal.length']) for r in rs]
    ys = [float(r['petal.width']) for r in rs]
    ax.scatter(xs, ys, label=v, color=COLORS[v], alpha=0.8, edgecolors='white', s=50)
ax.set_xlabel('花弁の長さ (cm)', fontsize=12)
ax.set_ylabel('花弁の幅 (cm)', fontsize=12)
ax.set_title('花弁の長さ × 幅 の散布図(品種別)', fontsize=14)
ax.legend()
fig.tight_layout()
fig.savefig('chart_scatter.png', dpi=150)
plt.close(fig)

# ---------- グラフ3: 箱ひげ図(花弁の長さ) ----------
fig, ax = plt.subplots(figsize=(9, 4.8))
data = [[float(r['petal.length']) for r in by_variety[v]] for v in varieties]
bp = ax.boxplot(data, tick_labels=varieties, patch_artist=True)
for patch, v in zip(bp['boxes'], varieties):
    patch.set_facecolor(COLORS[v])
    patch.set_alpha(0.7)
ax.set_ylabel('花弁の長さ (cm)', fontsize=12)
ax.set_title('品種別 花弁の長さの分布(箱ひげ図)', fontsize=14)
fig.tight_layout()
fig.savefig('chart_box.png', dpi=150)
plt.close(fig)

# ================== PPT作成 ==================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GRAY = RGBColor(0x44, 0x44, 0x44)


def add_title_bar(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.25), Inches(12.3), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_bullets(slide, items, left=0.7, top=1.6, width=11.9, height=5.3, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        indent = item[0]
        text = item[1]
        p.text = ('• ' if indent == 0 else '  – ') + text
        p.font.size = Pt(size if indent == 0 else size - 2)
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(10)


# ---- スライド1: 表紙 ----
slide = prs.slides.add_slide(BLANK)
rect = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
rect.fill.solid()
rect.fill.fore_color.rgb = NAVY
rect.line.fill.background()
box = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(11.3), Inches(2))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'Iris(アヤメ)データセット'
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p2 = tf.add_paragraph()
p2.text = '概要と分析結果'
p2.font.size = Pt(26)
p2.font.color.rgb = RGBColor(0xE0, 0xE6, 0xF0)

# ---- スライド2: データセット概要 ----
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, 'データセット概要')
add_bullets(slide, [
    (0, 'アヤメ(Iris)の花に関する古典的なデータセット(全150件)'),
    (0, '3品種 × 各50件で構成: Setosa / Versicolor / Virginica'),
    (0, '4つの数値特徴量を収録'),
    (1, 'sepal.length(がく片の長さ) / sepal.width(がく片の幅)'),
    (1, 'petal.length(花弁の長さ) / petal.width(花弁の幅)'),
    (0, '欠損値なし、単位は cm'),
    (0, '機械学習の分類問題における入門的ベンチマークとして広く利用される'),
], size=20)

# ---- スライド3: 全体統計 ----
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '全体の統計サマリー', '4つの特徴量の基本統計量(全150件)')
rows_tbl = 1 + len(COLS)
cols_tbl = 5
table_shape = slide.shapes.add_table(rows_tbl, cols_tbl, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.6))
table = table_shape.table
headers = ['特徴量', '平均', '標準偏差', '最小', '最大']
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(16)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
for i, c in enumerate(COLS):
    mean, std, mn, mx = overall_stats[c]
    vals = [COL_JP[c], f'{mean:.2f}', f'{std:.2f}', f'{mn:.2f}', f'{mx:.2f}']
    for j, v in enumerate(vals):
        cell = table.cell(i + 1, j)
        cell.text = v
        cell.text_frame.paragraphs[0].font.size = Pt(15)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF5, 0xFA) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
note = slide.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.5))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '品種は Setosa / Versicolor / Virginica が各50件ずつで、クラスの偏りがない均衡データセット'
p.font.size = Pt(16)
p.font.color.rgb = GRAY

# ---- スライド4: 品種ごとの平均値(棒グラフ) ----
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '品種ごとの特徴量比較')
slide.shapes.add_picture('chart_bar.png', Inches(1.1), Inches(1.5), width=Inches(11.1))

# ---- スライド5: 散布図 ----
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '花弁サイズによる品種の分離')
slide.shapes.add_picture('chart_scatter.png', Inches(1.1), Inches(1.5), width=Inches(11.1))

# ---- スライド6: 箱ひげ図 ----
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '花弁の長さの分布')
slide.shapes.add_picture('chart_box.png', Inches(1.1), Inches(1.5), width=Inches(11.1))

# ---- スライド7: まとめ ----
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, 'まとめ・考察')
add_bullets(slide, [
    (0, '花弁(petal)のサイズは品種間で差が大きく、分類に有効な特徴量'),
    (1, 'Setosa は花弁が明確に小さく、他2品種と容易に分離できる'),
    (1, 'Versicolor と Virginica はやや重なりがあるが、花弁の長さ・幅で概ね分離可能'),
    (0, 'がく片(sepal)のサイズは品種間の差が比較的小さく、識別力は花弁より弱い'),
    (0, 'データは3品種×50件で均衡しており、分類モデルの学習・評価に適している'),
    (0, '次のステップ: ロジスティック回帰やSVM等による品種分類モデルの構築'),
], size=20)

prs.save('iris_analysis.pptx')
print('Saved iris_analysis.pptx')
